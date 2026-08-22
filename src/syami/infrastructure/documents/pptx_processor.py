from pathlib import Path

from pptx import Presentation

from syami.application.documents.processor import DocumentProcessor
from syami.domain.document import ContentUnit, ProcessedDocument


class PPTXProcessor(DocumentProcessor):

    def process(
        self,
        document_id: int,
        source_path: str,
    ) -> ProcessedDocument:

        path = Path(source_path)

        if not path.exists():
            raise FileNotFoundError(
                f"File does not exist: {path}"
            )

        if path.suffix.lower() != ".pptx":
            raise ValueError(
                f"Expected a PPTX file: {path}"
            )

        presentation = Presentation(str(path))

        units = []

        for slide_index, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            text_parts = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text = shape.text.strip()

                if text:
                    text_parts.append(text)

            slide_text = "\n".join(text_parts)

            if not slide_text:
                continue

            units.append(
                ContentUnit(
                    text=slide_text,
                    unit_type="slide",
                    unit_index=slide_index,
                )
            )

        return ProcessedDocument(
            document_id=document_id,
            source_path=str(path),
            document_type="pptx",
            title=path.stem,
            units=units,
        )